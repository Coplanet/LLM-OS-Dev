from pathlib import Path
from typing import IO, Any, List, Union

from phi.document.base import Document
from phi.document.reader.base import Reader
from phi.utils.log import logger
from pptx import Presentation


class PPTXReader(Reader):
    """Reader for Excel files"""

    def read(self, path: Union[str, Path, IO[Any]]) -> List[Document]:
        if not path:
            raise ValueError("No path provided")

        try:
            logger.info(f"Reading: {path}")
            ppt_name = path.name.split(".")[0]
            ppt_contents = Presentation(path)
            text_content = []

            # Iterate through slides
            for slide in ppt_contents.slides:
                slide_text = []
                # Iterate through shapes in each slide
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                if slide_text:
                    # Join the text from each shape and add to the list
                    text_content.append("\n".join(slide_text))

            documents = [
                Document(
                    name=ppt_name,
                    id=f"{ppt_name}_{page_number}",
                    meta_data={"page": page_number},
                    content=content,
                )
                for page_number, content in enumerate(text_content, start=1)
            ]
            if self.chunk:
                chunked_documents = []
                for document in documents:
                    chunked_documents.extend(self.chunk_document(document))
                return chunked_documents

            return documents
        except Exception:
            raise
